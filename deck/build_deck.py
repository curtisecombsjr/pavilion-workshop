#!/usr/bin/env python3
"""
Build the Pavilion workshop deck as a .pptx (upload to Google Slides).

Design goals: SIMPLE and Google-Slides-safe — plain layouts, standard fonts
(Arial / Consolas), no fancy positioning. Content lives in the SLIDES list below,
so editing/adding/reordering is easy; re-run to regenerate.

    ./.venv/bin/python build_deck.py        # -> pavilion-workshop.pptx

Slide helpers: title / section / bullets / code / demo. Demo slides render a
terminal TRANSCRIPT — each command ($ prompt) followed by its real output —
from a `steps` list. Every slide can carry speaker notes (the `notes` key).
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------------- theme
NAVY   = RGBColor(0x1F, 0x3B, 0x63)
ACCENT = RGBColor(0x2E, 0x6D, 0xA4)
INK    = RGBColor(0x21, 0x25, 0x2B)
MUTED  = RGBColor(0x5B, 0x66, 0x74)
CODEBG = RGBColor(0xF3, 0xF4, 0xF6)
CODEINK= RGBColor(0x1B, 0x2B, 0x34)
GREEN  = RGBColor(0x1E, 0x7D, 0x3C)
RED    = RGBColor(0xB3, 0x2A, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD5, 0xDA, 0xE0)

BODY_FONT = "Arial"
CODE_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(run, text, size, color=INK, bold=False, font=BODY_FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _titlebar(slide, title, kicker=None):
    if kicker:
        _, tf = _box(slide, Inches(0.6), Inches(0.35), Inches(8.4), Inches(0.4))
        _set(tf.paragraphs[0].add_run(), kicker.upper(), 12, ACCENT, bold=True)
    _, tf = _box(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(0.9))
    _set(tf.paragraphs[0].add_run(), title, 30, NAVY, bold=True)
    ln = slide.shapes.add_shape(1, Inches(0.62), Inches(1.55), Inches(2.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def title_slide(spec):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    _, tf = _box(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2))
    _set(tf.paragraphs[0].add_run(), spec["title"], 46, WHITE, bold=True)
    p = tf.add_paragraph(); _set(p.add_run(), spec.get("subtitle", ""), 22, RGBColor(0xBF, 0xD3, 0xEA))
    links = spec.get("links")
    if links:
        _, tfl = _box(s, Inches(0.95), Inches(3.75), Inches(11.5), Inches(0.5))
        pl = tfl.paragraphs[0]
        for i, (text, url) in enumerate(links):
            if i:
                _set(pl.add_run(), "    ·    ", 16, RGBColor(0x6F, 0x86, 0xA6))
            # Plain text (no active hyperlink) so our light color sticks instead of the app's link-blue.
            _set(pl.add_run(), text, 16, RGBColor(0x9F, 0xD0, 0xFF))
    _, tf2 = _box(s, Inches(0.95), Inches(6.3), Inches(11), Inches(0.6))
    _set(tf2.paragraphs[0].add_run(), spec.get("footer", ""), 14, RGBColor(0x9F, 0xB6, 0xD2))
    _notes(s, spec.get("notes"))


def section_slide(spec):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.9), SW, Inches(1.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF7); bar.line.fill.background()
    _, tf = _box(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.3))
    _set(tf.paragraphs[0].add_run(), spec.get("kicker", "SECTION").upper(), 13, ACCENT, bold=True)
    p = tf.add_paragraph(); _set(p.add_run(), spec["title"], 36, NAVY, bold=True)
    _notes(s, spec.get("notes"))


def bullets_slide(spec):
    s = prs.slides.add_slide(BLANK)
    _titlebar(s, spec["title"], spec.get("kicker"))
    _, tf = _box(s, Inches(0.7), Inches(1.85), Inches(12), Inches(5.2))
    for i, b in enumerate(spec["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        lvl = 0
        if isinstance(b, tuple):
            b, lvl = b
        run = p.add_run()
        bullet = ("    – " if lvl else "•  ") + b
        _set(run, bullet, 20 if lvl == 0 else 17, INK if lvl == 0 else MUTED, bold=False)
    _notes(s, spec.get("notes"))


def code_slide(spec):
    s = prs.slides.add_slide(BLANK)
    _titlebar(s, spec["title"], spec.get("kicker"))
    if spec.get("intro"):
        _, tf = _box(s, Inches(0.7), Inches(1.75), Inches(12), Inches(0.5))
        _set(tf.paragraphs[0].add_run(), spec["intro"], 16, MUTED)
    top = Inches(2.35) if spec.get("intro") else Inches(1.95)
    h = spec.get("code_h", 4.6)
    box = s.shapes.add_shape(1, Inches(0.7), top, Inches(12), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
    for i, ln in enumerate(spec["code"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        color = CODEINK
        if ln.strip().startswith("#"):
            color = MUTED
        _set(p.add_run(), ln if ln else " ", spec.get("code_size", 14), color, font=CODE_FONT)
    _notes(s, spec.get("notes"))


def demo_slide(spec):
    """A live-demo slide as a terminal transcript: each command ($ prompt) followed by its real
    output, interleaved in one monospace box (kept sized to fit — no overflow)."""
    s = prs.slides.add_slide(BLANK)
    _titlebar(s, spec["title"], "LIVE DEMO · " + spec.get("level", ""))
    # run-cue chip (top-right) — only when there's a script to cue
    if spec.get("run"):
        chip = s.shapes.add_shape(1, Inches(9.3), Inches(0.5), Inches(3.4), Inches(0.55))
        chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT; chip.line.fill.background()
        ctf = chip.text_frame; ctf.word_wrap = True
        _set(ctf.paragraphs[0].add_run(), "▶  " + spec["run"], 13, WHITE, bold=True)
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # transcript box
    box = s.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(12), Inches(spec.get("box_h", 4.75)))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.16)
    # flatten steps -> (text, kind)
    lines = []
    for i, step in enumerate(spec["steps"]):
        if i > 0:
            lines.append((" ", "blank"))
        lines.append(("$ " + step["cmd"], "cmd"))
        for ln in step.get("out", "").split("\n"):
            lines.append((ln if ln else " ", "out"))
    for i, (text, kind) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if kind == "cmd":
            _set(p.add_run(), text, spec.get("cmd_size", 14), ACCENT, bold=True, font=CODE_FONT)
        else:
            color = CODEINK
            if "PASS" in text:
                color = GREEN
            elif "FAIL" in text:
                color = RED
            _set(p.add_run(), text, spec.get("out_size", 12.5), color, font=CODE_FONT)
    if spec.get("caption"):
        _, tf2 = _box(s, Inches(0.7), Inches(6.72), Inches(12), Inches(0.55))
        _set(tf2.paragraphs[0].add_run(), spec["caption"], 14, MUTED)
    _notes(s, spec.get("notes"))


RENDER = {"title": title_slide, "section": section_slide, "bullets": bullets_slide,
          "code": code_slide, "demo": demo_slide}

# ============================================================================ CONTENT
SLIDES = [
 {"type": "title", "title": "Pavilion",
  "subtitle": "A YAML-driven test framework for HPC systems",
  "links": [("github.com/hpc/pavilion2", "https://github.com/hpc/pavilion2"),
            ("pavilion2.readthedocs.io", "https://pavilion2.readthedocs.io")],
  "footer": "Workshop · live demos on the PBS cluster",
  "notes": "Set expectations: ~60 min, concept + lots of live demos. We'll go from a "
           "one-line 'hello' test all the way to custom plugins shipping results to OpenSearch."},

 {"type": "section", "title": "Why Pavilion?", "kicker": "The problem"},
 {"type": "bullets", "title": "What Pavilion is", "kicker": "Overview",
  "bullets": [
    "An open-source (LANL) framework for testing HPC systems.",
    "Mature & battle-tested — years of development at LANL, hardened on production supercomputers.",
    "Tests are YAML configs that wrap your test codes — not scripts you maintain by hand.",
    "Runs the same tests across systems and schedulers (Slurm, Flux, PBS).",
    "Plugin-driven: schedulers, result parsers, result loggers, commands, system variables.",
    "Inheritable test definitions — write shared config once; tests inherit and override it (DRY).",
  ],
  "notes": "The pitch for a mixed audience: stop hand-writing brittle test scripts. Describe "
           "a test once in YAML; Pavilion handles building, scheduling, running, parsing, and "
           "recording results — consistently, everywhere."},

 {"type": "bullets", "title": "Built for", "kicker": "Where it fits",
  "bullets": [
    "Acceptance testing — does a new or returned system meet spec before users get it?",
    "System validation — confirm the whole stack still behaves after maintenance.",
    "Baseline testing — capture known-good numbers to measure future runs against.",
    "Regression testing — catch when a change quietly breaks or slows something.",
    "Reproducibility — the same test, run the same way, by anyone, every time.",
  ],
  "notes": "The 'why should I care' slide. Baseline testing pairs with regression: you can't "
           "spot a regression without a known-good baseline to compare to."},

 {"type": "bullets", "title": "The problem it solves", "kicker": "Why it matters",
  "bullets": [
    "A new cluster, a maintenance window, or a hardware repair — does everything still work and perform?",
    "Hand-rolled scripts drift, differ per person, and don't record results.",
    "Pavilion records every run — a searchable history you filter by result, name, or date.",
    "Pavilion gives one declarative definition + a full run/results pipeline.",
    "Results are captured, parsed, and shippable to dashboards (OpenSearch → Grafana).",
  ]},

 {"type": "section", "title": "Core concepts", "kicker": "How it fits together"},
 {"type": "bullets", "title": "The building blocks", "kicker": "Concepts",
  "bullets": [
    "Test — one thing to run, defined in YAML (build + run + result parsing).",
    "Suite — a YAML file holding related tests (e.g. demo_pbs.yaml).",
    "Config layering — host → mode → test: defaults get overridden in that order.",
    "Variables & permutations — one definition expands into many test instances.",
    "config_dirs — where Pavilion looks for suites, modes, series, and plugins.",
  ],
  "notes": "Name each concept plainly. 'suite.test' is how you address a test, e.g. "
           "demo_pbs.pass."},

 {"type": "bullets", "title": "The test lifecycle", "kicker": "Under the hood",
  "bullets": [
    "1. Resolve — read configs, apply host/mode overrides, expand permutations.",
    "2. Build — set up the test's source/build once (reused across runs).",
    "3. Kickoff — hand the test to the scheduler (raw / PBS / Slurm / Flux).",
    "4. Run — execute run.cmds on the target node(s).",
    "5. Results — parse output → results.json, then hand off to result loggers.",
  ],
  "notes": "This is the mental model for the demos: run -> build (reused) -> kickoff -> run -> "
           "parse -> log."},

 {"type": "bullets", "title": "The build phase", "kicker": "Under the hood",
  "bullets": [
    "Before a test runs, Pavilion builds it — once — then reuses that build across runs.",
    "Compile from source — hand it a source archive/URL + build commands; it builds for you.",
    "Build from a source directory — point at a local dir and run your build steps there.",
    "Or just stage binaries — copy in a directory of pre-built binaries, no compile at all.",
    "Builds are cached & shared — reruns skip rebuilding (the 'BUILD_REUSED' you'll see live).",
  ],
  "notes": "The build step is distinct from the run. Three flavors: compile from source, build "
           "from a source directory, or copy pre-built binaries. Built once (on the login node — "
           "compute nodes often have no internet) and reused, which is why demos show BUILD_REUSED."},

 {"type": "code", "title": "Anatomy of a test (YAML)", "kicker": "Concepts",
  "intro": "A trivial test: a name, a scheduler + queue, and the commands to run.",
  "code":
    "demo_pass:\n"
    "  summary: 'Trivial PASS via PBS'\n"
    "  scheduler: pbs\n"
    "  schedule:\n"
    "    pbs:\n"
    "      queue: workq          # submit to a PBS queue\n"
    "  run:\n"
    "    cmds:\n"
    "      - 'echo \"hello from $(hostname)\"'\n"
    "\n"
    "# PASS/FAIL comes from the exit code of the last command.\n"
    "# Address it as:  suite.test   ->   demo_pbs.pass",
  "notes": "Emphasize: the exit code decides PASS/FAIL. It's dispatched to PBS (queue workq)."},

 {"type": "code", "title": "A fuller test (YAML)", "kicker": "Concepts",
  "intro": "Same shape, more of it — variables, PBS + a queue, and build → run → parse → evaluate.",
  "code_size": 13.5, "code_h": 5.0,
  "code":
    "stream:\n"
    "  variables:\n"
    "    array_size: 40000000\n"
    "    threads: 4\n"
    "  scheduler: pbs\n"
    "  schedule:\n"
    "    pbs:\n"
    "      queue: workq          # target a queue\n"
    "      nodes: 1\n"
    "      walltime: '00:05:00'\n"
    "  build:                    # compile from source\n"
    "    source_path: stream.c\n"
    "    cmds: ['gcc -fopenmp -DN={{array_size}} stream.c -o stream']\n"
    "  run:\n"
    "    env: {OMP_NUM_THREADS: '{{threads}}'}\n"
    "    cmds: ['./stream']\n"
    "  result_parse:\n"
    "    regex:\n"
    "      triad_mbs: {regex: 'Triad:\\s+(\\S+)'}\n"
    "  result_evaluate:\n"
    "    result: 'triad_mbs > 10000'    # PASS only if fast enough",
  "notes": "The fuller shape: a couple of variables up top (used below), the PBS scheduler targeting "
           "a queue, then the four stages — build (compile from source), run, result_parse (regex → "
           "triad_mbs), and result_evaluate (a pass/fail expression). Still just YAML."},

 {"type": "bullets", "title": "Inheritable test definitions", "kicker": "A big one",
  "bullets": [
    "Write shared config once in a base test — scheduler, build, run, parsing.",
    "Other tests just say  inherits_from: <base>  and get all of it for free.",
    "Each one overrides only what changes — no copy-paste, no drift.",
    "Fix the base once → every test that inherits it updates.",
  ],
  "notes": "Huge selling point. Real suites share tons of config (same build, same scheduler "
           "settings). Inheritance = write it once, variants override just the delta. Less code, "
           "less drift, one place to fix things."},

 {"type": "code", "title": "Inheritance in action", "kicker": "A big one",
  "code_size": 12.5, "code_h": 5.45,
  "code":
    "base:                        # shared config, written once\n"
    "  variables:\n"
    "    threads: 2\n"
    "  build:\n"
    "    cmds: ['gcc -O3 -fopenmp stream.c -o stream']\n"
    "  run:\n"
    "    cmds: ['OMP_NUM_THREADS={{threads}} ./stream']\n"
    "\n"
    "small:\n"
    "  inherits_from: base        # that's it — inherits everything\n"
    "\n"
    "large:\n"
    "  inherits_from: base\n"
    "  variables:\n"
    "    threads: 5               # override just the variable\n"
    "\n"
    "debug:\n"
    "  inherits_from: base\n"
    "  build:                     # override just the build flags\n"
    "    cmds: ['gcc -O3 -fopenmp -fsanitize=address stream.c -o stream']",
  "notes": "One base owns the shared config and a variable (threads: 2). 'small' is literally just "
           "inherits_from: base — it takes everything. 'large' inherits and overrides only the variable "
           "(threads: 5). 'debug' inherits and overrides only the build, adding -fsanitize=address. "
           "Each child writes just its delta — that's the DRY win."},

 {"type": "bullets", "title": "The CLI you'll use", "kicker": "Concepts",
  "bullets": [
    "pav run <suite.test>   — build + run a test (or a whole suite).",
    "pav status <sid>       — where are my tests / did they pass?",
    "pav results <sid>      — the results table; --full for the full JSON.",
    "pav wait <sid>         — block until tests finish.",
    "pav show tests|sched|modes|series   — what's available.",
    ("pav log, pav cat, pav ls, pav series … — inspect a run in detail.", 1),
  ],
  "notes": "Tip: 'pav --quiet' silences a harmless config-label warning for clean output."},

 {"type": "demo", "title": "Your whole run history — searchable", "level": "RECORDS",
  "run": "11-history.sh",
  "steps": [
    {"cmd": "pav status all -F failed                    # every failed run, all the way back",
     "out": "s70.2   FAIL   demo_pbs.fail              16:40\n"
            "s75.2   FAIL   opensearch_verify.os_fail  16:44\n"
            "s23.71  FAIL   sim_tests.sim_fail_01      Apr 18"},
    {"cmd": "pav status all -F 'name=demo_pbs.metrics'   # one test, every run",
     "out": "s71.1   PASS   demo_pbs.metrics    16:40\n"
            "s67.2   PASS   demo_pbs.metrics    16:32\n"
            "s57.1   PASS   demo_pbs.metrics    16:28"},
    {"cmd": "pav status all -F 'name=demo_pbs.fail and failed'   # combine filters",
     "out": "s70.2   FAIL   demo_pbs.fail    16:40\n"
            "s56.2   FAIL   demo_pbs.fail    16:27"},
  ],
  "caption": "Every run is kept. 'all' searches the full history; -F filters — and combines — by result, name, date…",
  "notes": "Pavilion records every test run persistently. 'pav status all' searches the whole history "
           "(not just the last day); -F takes keywords (PASSED/FAILED) or field filters (name=…), joined "
           "with 'and'/'or'. Note the Apr-26 failure — the record goes back as far as your runs do."},

 {"type": "demo", "title": "Basics: PASS and FAIL", "level": "L1", "run": "01-basic.sh",
  "steps": [
    {"cmd": "pav run demo_pbs.pass demo_pbs.fail          # two PBS jobs",
     "out": "sid: s70\ntests: 2"},
    {"cmd": "qstat -a                                     # yes — it went through the scheduler",
     "out": "Job ID          Username Queue Jobname    NDS TSK S\n"
            "8383.pbs-server pavilion workq pav_demo_*   1   1 R"},
    {"cmd": "pav status s70",
     "out": "s70.1 | 8383_pbs-server | demo_pbs.pass | COMPLETE | PASS\n"
            "s70.2 | 8383_pbs-server | demo_pbs.fail | COMPLETE | FAIL"},
  ],
  "caption": "Submit two PBS jobs, prove dispatch with qstat, read the results. FAIL is intentional (exit 7).",
  "notes": "Two tests from one 'pav run', both submitted to PBS (qstat proves it). PASS vs FAIL is "
           "purely the exit code. Running the PBS test also proves the scheduler itself works."},

 {"type": "section", "title": "Getting data out", "kicker": "Result parsing"},
 {"type": "code", "title": "Result parsing", "kicker": "Results",
  "intro": "Regex parsers pull values out of the test's output into structured results.",
  "code":
    "metrics:\n"
    "  scheduler: pbs\n"
    "  schedule:\n"
    "    pbs: { queue: workq, nodes: 1, walltime: '00:02:00' }\n"
    "  run:\n"
    "    cmds:\n"
    "      - 'echo \"throughput 123.45\"'\n"
    "      - 'echo \"latency 5.6\"'\n"
    "  result_parse:\n"
    "    regex:\n"
    "      throughput_mbs: { regex: 'throughput (\\S+)', action: store }\n"
    "      latency_ms:     { regex: 'latency (\\S+)',    action: store }",
  "notes": "Many parsers exist (regex, constant, table, split…) plus result_evaluate for computed "
           "pass/fail. Raw stdout becomes queryable numbers."},

 {"type": "demo", "title": "Parsed metrics in results", "level": "L2", "run": "02-metrics.sh",
  "steps": [
    {"cmd": "pav run demo_pbs.metrics",
     "out": "sid: s71\ntests: 1"},
    {"cmd": "qstat -a                         # dispatched to PBS",
     "out": "8384.pbs-server pavilion workq pav_demo_*  1  1  R"},
    {"cmd": "pav results --full s71",
     "out": "  'name': 'demo_pbs.metrics',\n"
            "  'result': 'PASS',\n"
            "  'throughput_mbs': 123.45,\n"
            "  'latency_ms': 5.6,\n"
            "  'errors': 0,"},
  ],
  "caption": "Those numbers are now in results.json — ready for logging and dashboards.",
  "notes": "These parsed fields are exactly what ships to OpenSearch later."},

 {"type": "section", "title": "One config → many tests", "kicker": "Permutations"},
 {"type": "code", "title": "Permutations", "kicker": "The big one",
  "intro": "permute_on + variable lists generate one test instance per combination.",
  "code":
    "matrix:\n"
    "  scheduler: pbs\n"
    "  schedule:\n"
    "    share_allocation: false      # one PBS job per instance\n"
    "    pbs: { queue: workq, nodes: 1, walltime: '00:02:00' }\n"
    "  permute_on: [size, mode]\n"
    "  variables:\n"
    "    size: ['small', 'medium', 'large']\n"
    "    mode: ['read', 'write']\n"
    "  subtitle: '{{size}}-{{mode}}'\n"
    "  run:\n"
    "    cmds:\n"
    "      - 'echo \"size={{size}} mode={{mode}}\"'\n"
    "\n"
    "# 3 sizes x 2 modes  =>  6 test instances, 6 PBS jobs.",
  "notes": "The feature people remember. One definition, six tests, each with its own subtitle. "
           "share_allocation: false makes Pavilion submit a separate PBS job per instance (the default "
           "packs them into one shared allocation)."},

 {"type": "demo", "title": "Permutations expand", "level": "L3", "run": "03-permutations.sh",
  "steps": [
    {"cmd": "pav run demo_perms.matrix",
     "out": "Test set 'demo_perms.matrix' created 6 tests, skipped 0, 0 errors.\nsid: s93"},
    {"cmd": "qstat -a                        # SIX PBS jobs — one per permutation",
     "out": "8400.pbs-server pavilion workq pav_demo_*  1  1  R\n"
            "8401.pbs-server pavilion workq pav_demo_*  1  1  R\n"
            "8402.pbs-server pavilion workq pav_demo_*  1  1  Q\n"
            "8403.pbs-server pavilion workq pav_demo_*  1  1  Q\n"
            "8404.pbs-server pavilion workq pav_demo_*  1  1  Q\n"
            "8405.pbs-server pavilion workq pav_demo_*  1  1  Q"},
    {"cmd": "pav status s93                  # each instance = its own job id",
     "out": "s93.1 | 8400_pbs-server | demo_perms.matrix.small-read  | PASS\n"
            "s93.2 | 8401_pbs-server | demo_perms.matrix.medium-read | PASS\n"
            "s93.3 | 8402_pbs-server | demo_perms.matrix.large-read  | PASS\n"
            "s93.4 | 8403_pbs-server | demo_perms.matrix.small-write | PASS\n"
            "s93.5 | 8404_pbs-server | demo_perms.matrix.medium-write| PASS\n"
            "s93.6 | 8405_pbs-server | demo_perms.matrix.large-write | PASS"},
  ],
  "caption": "One YAML block → six PBS jobs (8400–8405), one per permutation — some run, the rest queue.",
  "notes": "share_allocation: false gives each permutation its OWN PBS job (distinct job ids 8400–8405). "
           "With one free node they run a couple at a time and the rest sit in Q — real scheduler behavior. "
           "Point out the auto-generated subtitles (size-mode)."},

 {"type": "section", "title": "Schedulers", "kicker": "Running at scale"},
 {"type": "bullets", "title": "Schedulers & the PBS plugin", "kicker": "Schedulers",
  "bullets": [
    "Built-in scheduler plugins: raw (local), slurm, flux, pbs.",
    "Switch by setting `scheduler:` — the test body doesn't change.",
    "The PBS plugin (ours) submits real qsub jobs and maps nodes/tasks/walltime/queue.",
    "Scheduler variables (sched.*) expose node lists, chunks, task counts to the test.",
  ],
  "notes": "Same test, different scheduler = portability. Our PBS plugin is the custom bit."},

 {"type": "code", "title": "A PBS test", "kicker": "Schedulers",
  "code":
    "# demo_pbs.yaml  ->  address it as  demo_pbs.pass\n"
    "pass:\n"
    "  scheduler: pbs\n"
    "  schedule:\n"
    "    pbs:\n"
    "      nodes: 1\n"
    "      tasks: 1\n"
    "      walltime: '00:00:30'\n"
    "      queue: workq\n"
    "  run:\n"
    "    cmds:\n"
    "      - 'echo \"pbs job on $(hostname)\"'",
  "notes": "schedule.pbs maps to #PBS resource requests via our plugin. Addressed as demo_pbs.pass."},

 {"type": "demo", "title": "Submit to PBS", "level": "L4", "run": "04-pbs.sh",
  "steps": [
    {"cmd": "pav run demo_pbs.pass",
     "out": "sid: s76\ntests: 1"},
    {"cmd": "pav status s76           # Pavilion: it's running",
     "out": "s76.1 | 8390_pbs-server | demo_pbs.pass | RUNNING"},
    {"cmd": "qstat -a                 # PBS: the same job, state R",
     "out": "8390.pbs-server  pavilion  workq  pav_demo_*  1  1  R"},
    {"cmd": "pav results s76          # once it finishes",
     "out": "s76.1 | demo_pbs.pass | COMPLETE | PASS"},
  ],
  "caption": "Submitted → running (seen in Pavilion AND PBS) → PASS on a compute node.",
  "notes": "The running step shows the same job in both Pavilion (RUNNING) and PBS (state R). "
           "Live tip: 'pav run' then immediately 'pav status' / 'qstat -a' to catch it."},

 {"type": "section", "title": "Modes", "kicker": "Reusable overlays"},
 {"type": "bullets", "title": "What a mode is", "kicker": "Modes",
  "bullets": [
    "A mode is a small YAML overlay merged on top of a fully-resolved test.",
    "Apply at run time with -m: `pav run -m prod <test>`.",
    "Great for: swapping queues, tuning resources, toggling variables — without editing tests.",
    "Our demo: the test fails at throughput 1200; a `prod` mode raises it to 1400 → it passes.",
  ],
  "notes": "Modes = reuse. One test, many contexts. Here the same test flips FAIL→PASS purely "
           "because the mode overrode a value — no edit to the test itself."},

 {"type": "code", "title": "The test, and a mode that changes its outcome", "kicker": "Modes",
  "intro": "The test measures throughput and requires > 1300. Base is 1200 (FAIL); the mode raises it.",
  "code":
    "# config/suites/demo_modes.yaml   (the test)\n"
    "mode_demo:\n"
    "  variables:\n"
    "    throughput: 1200\n"
    "  run:\n"
    "    cmds:\n"
    "      - 'echo \"throughput {{throughput}}\"'\n"
    "  result_parse:\n"
    "    regex:\n"
    "      throughput_mbs:\n"
    "        regex: 'throughput (\\S+)'\n"
    "  result_evaluate:\n"
    "    result: 'throughput_mbs > 1300'      # PASS only if above 1300\n"
    "\n"
    "# config/modes/prod.yaml   (the mode overlay)\n"
    "variables:\n"
    "  throughput: 1400                        # override 1200 -> 1400",
  "code_size": 13, "code_h": 5.3,
  "notes": "The test evaluates throughput_mbs > 1300. Base value 1200 fails it. The prod mode "
           "overrides only that variable to 1400, so the same test now passes — no edit to the test."},

 {"type": "demo", "title": "A mode flips FAIL → PASS", "level": "L5", "run": "05-modes.sh",
  "steps": [
    {"cmd": "pav run demo_modes.mode_demo           # base: throughput = 1200",
     "out": "  'throughput_mbs': 1200,\n"
            "  'result': 'FAIL'          # 1200 > 1300 is false"},
    {"cmd": "qstat -a                                # a real PBS job",
     "out": "8407.pbs-server pavilion workq pav_demo_*  1  1  R"},
    {"cmd": "pav run -m prod demo_modes.mode_demo    # prod overrides -> throughput = 1400",
     "out": "  'throughput_mbs': 1400,\n"
            "  'result': 'PASS'          # 1400 > 1300 is true"},
  ],
  "caption": "Same test, no edits — the prod mode raised throughput past the 1300 threshold, flipping FAIL → PASS.",
  "notes": "The result is the proof: base throughput 1200 fails the > 1300 check; the prod mode overrides "
           "the value to 1400 and the identical test now passes. Both ran through PBS (qstat)."},

 {"type": "section", "title": "Series", "kicker": "Grouping runs"},
 {"type": "code", "title": "A series file", "kicker": "Series",
  "intro": "demo_series.yaml — name a few test sets; each set lists the tests it runs.",
  "code":
    "# config/series/demo_series.yaml\n"
    "ordered: False\n"
    "test_sets:\n"
    "  smoke:\n"
    "    tests:\n"
    "      - demo_pbs.pass\n"
    "  perf:\n"
    "    tests:\n"
    "      - demo_pbs.metrics",
  "notes": "Sets can be ordered or not; unordered runs them concurrently (robust for a live demo). "
           "Each set becomes its own PBS job."},

 {"type": "demo", "title": "Run a whole group", "level": "L6", "run": "06-series.sh",
  "steps": [
    {"cmd": "pav series run demo_series",
     "out": "Created Test Series demo_series.\nStarted series s74."},
    {"cmd": "qstat -a                        # one PBS job per test set",
     "out": "8387.pbs-server pavilion workq pav_demo_*  1  1  R\n"
            "8388.pbs-server pavilion workq pav_demo_*  1  1  R"},
    {"cmd": "pav series status s74",
     "out": "Id  | Name        | Status   | Tests | Pass | Fail\n"
            "s74 | demo_series | COMPLETE |   2   |   2  |   0"},
  ],
  "caption": "One command runs both test sets as a named series — two PBS jobs, all PASS.",
  "notes": "Series group test sets; each set here dispatches its own PBS job (qstat shows two). "
           "I use unordered for a robust live run."},

 {"type": "section", "title": "Results", "kicker": "Reading & scripting"},
 {"type": "demo", "title": "Results: human and machine", "level": "RESULTS", "run": None,
  "out_size": 11,
  "steps": [
    {"cmd": "pav results s71                     # a clean table for humans",
     "out": " Test Results: s71.\n"
            " Id    | Name             | Started  | Result\n"
            " s71.1 | demo_pbs.metrics | 16:40:44 | PASS"},
    {"cmd": "pav results --json s71              # the same run, for scripting",
     "out": "{\n"
            "  \"name\": \"demo_pbs.metrics\",\n"
            "  \"result\": \"PASS\",\n"
            "  \"throughput_mbs\": 123.45,\n"
            "  \"latency_ms\": 5.6,\n"
            "  \"errors\": 0,\n"
            "  \"duration\": 0.041\n"
            "}"},
  ],
  "caption": "Same results, two shapes: a table to read, JSON to pipe into jq / a script / a logger.",
  "notes": "Point out that every parsed field (throughput_mbs, latency_ms…) is in the JSON — that's what "
           "feeds the result loggers next. --json (or -j) is how you script against Pavilion."},

 {"type": "section", "title": "Plugins", "kicker": "Extending Pavilion"},
 {"type": "bullets", "title": "Pavilion is plugins all the way down", "kicker": "Plugins",
  "bullets": [
    "Scheduler plugins — raw, slurm, flux, pbs (PBS is ours).",
    "Result parsers — regex, constant, table, …",
    "System variables — expose machine facts to tests (e.g. sys_name).",
    "Command plugins — add whole new `pav <command>` subcommands.",
    "Result loggers — ship results anywhere (CSV, OpenSearch, custom).",
    ("Drop a .py + a .yapsy-plugin into config/plugins/ — Pavilion discovers it.", 1),
  ],
  "notes": "Two kinds of custom plugins next: command plugins and result loggers."},

 {"type": "bullets", "title": "Command plugins (ours)", "kicker": "Custom commands",
  "bullets": [
    "hello — a friendly sanity check + Pavilion info.",
    "recent — the most recent test runs, colored by result.",
    "test-summary — a PASS/FAIL tally across recent runs.",
    "disk-usage — how much space test runs / builds / series consume.",
  ],
  "notes": "These add new verbs to pav itself — just Python classes Pavilion auto-discovers."},

 {"type": "demo", "title": "Custom commands in action", "level": "L7", "run": "07-command-plugins.sh",
  "out_size": 11.5,
  "steps": [
    {"cmd": "pav hello --name Team",
     "out": "Hello, Team!  Welcome to Pavilion 2!"},
    {"cmd": "pav recent -n 3",
     "out": "[s75.4] opensearch_verify.os_version_tag.v1.0-verify   PASS\n"
            "[s75.3] opensearch_verify.os_with_metrics             PASS\n"
            "[s75.2] opensearch_verify.os_fail                     FAIL"},
    {"cmd": "pav test-summary",
     "out": "PASS: 74    FAIL: 8    TOTAL: 83"},
    {"cmd": "pav disk-usage        # where the space goes",
     "out": "Config Area: main\n"
            "  Test Runs:     4.42 MB\n"
            "  Builds:       59.77 MB\n"
            "Series:          16.84 MB\n"
            "------------------------------------------------\n"
            "TOTAL USAGE:    145.22 MB"},
  ],
  "caption": "Four new pav subcommands — our own tooling built on Pavilion's plugin API.",
  "notes": "hello/recent/test-summary/disk-usage all add new verbs to pav itself. recent/test-summary "
           "read Pavilion's run database via cmd_utils; disk-usage totals runs / builds / series (handy "
           "before a 'pav clean'). Same plugin API for all four."},

 {"type": "bullets", "title": "Result loggers (ours)", "kicker": "Custom output",
  "bullets": [
    "csv_file — appends one CSV row per test (result, duration, metrics, permute vars).",
    "saLog — a custom NASA/NOAA logger; writes a CSV row on result 'released'.",
    "opensearch — ships each result document into the `pavilion-results` index.",
    "Configured once in pavilion.yaml → every test result flows through them.",
  ],
  "notes": "Result loggers run at the end of the results stage. One config, applies to all runs. "
           "(Redact the OpenSearch password if you ever show pavilion.yaml.)"},

 {"type": "demo", "title": "Output plugin: CSV", "level": "L8", "run": "08-output-csv.sh",
  "out_size": 11.5,
  "steps": [
    {"cmd": "pav run demo_pbs.pass demo_pbs.metrics      # real PBS jobs",
     "out": "sid: s74\ntests: 2"},
    {"cmd": "qstat -a                                    # yes — it went through the scheduler",
     "out": "8387.pbs-server  pavilion  workq  pav_demo_*  1  1  R"},
    {"cmd": "tail ~/pav_logs/results.csv                 # csv_file logger: one row per test",
     "out": "name,id,result,sys_name,user, ... ,duration\n"
            "demo_pbs.pass,s74.1,PASS,pbs-server,pavilion, ... ,0.043\n"
            "demo_pbs.metrics,s74.2,PASS,pbs-server,pavilion, ... ,0.040"},
  ],
  "caption": "A real PBS job (qstat proves dispatch), then the csv_file logger writes one row per test.",
  "notes": "The CSV page also shows scheduler dispatch — these are real PBS jobs (qstat) — then the "
           "csv_file logger records one row per test. Parsed metrics land in 'extra', permute vars in 'permute_on'."},

 {"type": "demo", "title": "Output plugin: saLog", "level": "saLog", "run": "12-salog.sh",
  "steps": [
    {"cmd": "pav run demo_pbs.salog                 # 1. run in Pavilion (per_file: name)",
     "out": "sid: s141\ntests: 1"},
    {"cmd": "qstat -a                               # 2. the job, running in PBS",
     "out": "8410.pbs-server pavilion workq pav_demo_*  1  1  R"},
    {"cmd": "tail salog_output.txt                  # 3. the CSV the sa_log plugin wrote on release",
     "out": "x1002c6s2b0n1,\"demo_pbs.salog result=PASS\",pavilion,released"},
  ],
  "caption": "Run in Pavilion → job in PBS → the sa_log plugin calls saLog per node as it's released.",
  "notes": "saLog is a logger specific to us (NASA/NOAA). The test uses a per_file: name parser so each "
           "node's result is recorded; Pavilion's sa_log plugin then invokes saLog once per node with the "
           "'released' action, writing node, content (result), user, action to the CSV."},

 {"type": "demo", "title": "Output plugin: OpenSearch → Grafana", "level": "L9", "run": "09-opensearch.sh",
  "steps": [
    {"cmd": "pav run -m pbs opensearch_verify        # 4 tests via PBS; each result → OpenSearch",
     "out": "created 4 tests, skipped 0, 0 errors.\nsid: s75"},
    {"cmd": "qstat -a                                # through the scheduler",
     "out": "8389.pbs-server pavilion workq pav_opens*  1  1  R"},
    {"cmd": "python3 ~/opensearch_results.py --name opensearch_verify",
     "out": "pav_id   name                                          result\n"
            "s75.1    opensearch_verify.os_pass                     PASS\n"
            "s75.2    opensearch_verify.os_fail                     FAIL\n"
            "s75.3    opensearch_verify.os_with_metrics             PASS\n"
            "s75.4    opensearch_verify.os_version_tag.v1.0-verify  PASS"},
  ],
  "caption": "Shipped to OpenSearch, read back from the index — then on to Grafana dashboards.",
  "notes": "Run the suite, query the index to prove indexing, then open Grafana "
           "(http://<your-grafana-host>:3000). The read-back auto-loads the password from "
           "pavilion.yaml — no OS_PASS export needed."},

 {"type": "demo", "title": "Output plugin: MySQL", "level": "L10", "run": "10-mysql.sh",
  "steps": [
    {"cmd": "pav run demo_pbs.metrics          # every logger fires, incl. mysql",
     "out": "sid: s71\ntests: 1"},
    {"cmd": "qstat -a                          # a real PBS job",
     "out": "8384.pbs-server pavilion workq pav_demo_*  1  1  R"},
    {"cmd": "mysql pavilion -e 'SELECT pav_id,name,result,sys_name,ROUND(duration,3) dur\n"
            "                    FROM results ORDER BY logged_at DESC LIMIT 4'",
     "out": "pav_id  name                  result  sys_name    dur\n"
            "s74.2   demo_pbs.metrics      PASS    pbs-server  0.040\n"
            "s74.1   demo_pbs.pass         PASS    pbs-server  0.043\n"
            "s73.1   demo_modes.mode_demo  PASS    pbs-server  0.040\n"
            "s72.3   demo_perms.matrix...  PASS    pbs-server  0.040"},
  ],
  "caption": "A result logger we wrote — every result lands in a MySQL table (passwordless local auth).",
  "notes": "Built with pymysql; connects via unix_socket so there's NO password in pavilion.yaml. "
           "The row-per-test appears the moment the run finishes."},

 {"type": "bullets", "title": "Export results anywhere", "kicker": "Result loggers",
  "bullets": [
    "Grafana — live dashboards & trends  (shown today).",
    "OpenSearch — index every result document  (shown today).",
    "MySQL — a row per result in a relational table  (shown today).",
    "PostgreSQL — same idea, for a Postgres shop.",
    "[add your own] — a result logger is just a plugin: target any destination.",
  ],
  "notes": "The point: loggers are plugins, so results can go wherever you need. csv_file / "
           "saLog / opensearch / mysql all exist and are shown today; PostgreSQL is an example "
           "you'd write a logger for — '[add your own]' says exactly that."},

 {"type": "section", "title": "Putting it together", "kicker": "Real-world"},
 {"type": "bullets", "title": "How we use it here", "kicker": "In practice",
  "bullets": [
    "Define acceptance/regression tests once (incl. real benchmarks: HPL, HPCG).",
    "Run through PBS across the cluster; permute across configurations.",
    "Parse metrics; every result flows to CSV + OpenSearch automatically.",
    "Grafana dashboards make trends and regressions visible over time.",
  ],
  "notes": "Pavilion isn't just running tests — it's a data pipeline from 'did it pass' to "
           "'how is the system trending'."},

 {"type": "bullets", "title": "Recap", "kicker": "Wrap-up",
  "bullets": [
    "YAML tests → build → run → parse → log. One model, everywhere.",
    "Permutations turn one config into a whole test matrix.",
    "Schedulers (incl. our PBS plugin) run it at cluster scale.",
    "Plugins extend everything — our commands and result loggers included.",
    "Results become dashboards. Questions?",
    "Learn more →  github.com/hpc/pavilion2   ·   pavilion2.readthedocs.io",
  ],
  "notes": "Leave up for Q&A. Offer to walk through any demo again live."},
]

for spec in SLIDES:
    RENDER[spec["type"]](spec)

# Save next to this script, so the build works from any working directory.
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pavilion-workshop.pptx")
prs.save(OUT)
print(f"Wrote {OUT} with {len(SLIDES)} slides.")
